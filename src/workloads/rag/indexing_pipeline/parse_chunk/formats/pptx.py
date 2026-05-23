#!/usr/bin/env python3
"""
pptx format parser — import-safe, lazy-init, compatible with router.py.

Key changes from the original:
- No sys.exit()/heavy work at import time.
- Environment, boto3, pyarrow, numpy are initialized lazily inside _init_env().
- parse_file() validates runtime requirements and raises on fatal problems.
- CLI behavior preserved under `if __name__ == "__main__":`
"""
from __future__ import annotations
import os
import sys
import json
import time
import logging
import hashlib
import tempfile
import re
import unicodedata
from io import BytesIO
from datetime import datetime
from typing import List, Dict, Any, Optional, Generator, Tuple
from botocore.exceptions import ClientError

# --- simple color logging (safe at import) ---------------------------------
RESET = "\033[0m"
COLORS = {
    logging.DEBUG: "\033[90m",
    logging.INFO: "\033[37m",
    logging.WARNING: "\033[33m",
    logging.ERROR: "\033[31m",
    logging.CRITICAL: "\033[1;41m",
}
class ColorFormatter(logging.Formatter):
    def format(self, record):
        color = COLORS.get(record.levelno, RESET)
        message = super().format(record)
        return f"{color}{message}{RESET}"

log = logging.getLogger("pptx_parser")
level_env = os.getenv("LOG_LEVEL", "INFO")
try:
    log.setLevel(level_env)
except Exception:
    log.setLevel("INFO")
handler = logging.StreamHandler(sys.stdout)
handler.setFormatter(ColorFormatter("%(asctime)s %(levelname)s %(message)s"))
log.handlers[:] = [handler]

# --- module-level defaults (will be set by _init_env) ----------------------
S3_BUCKET: Optional[str] = None
S3_RAW_PREFIX: str = os.getenv("S3_RAW_PREFIX", "data/raw/").rstrip("/") + "/"
S3_CHUNKED_PREFIX: str = os.getenv("S3_CHUNKED_PREFIX", "data/chunked/").rstrip("/") + "/"
SLIDES_PER_CHUNK: int = int(os.getenv("PPTX_SLIDES_PER_CHUNK", "3"))
DISABLE_OCR: bool = os.getenv("PPTX_DISABLE_OCR", "false").lower() == "true"
FORCE_OCR: bool = os.getenv("PPTX_FORCE_OCR", "false").lower() == "true"
OCR_BACKEND: str = os.getenv("PPTX_OCR_ENGINE", "tesseract").lower()
PPTX_OCR_STRICT: bool = os.getenv("PPTX_OCR_STRICT", "false").lower() == "true"
MIN_IMG_BYTES: int = int(os.getenv("PPTX_MIN_IMG_SIZE_BYTES", "3072"))
PARSER_VERSION_PPTX: str = os.getenv("PARSER_VERSION_PPTX", "pptx-parser-v1")
TOKEN_ENCODER: str = os.getenv("TOKEN_ENCODER", "cl100k_base")
FORCE_OVERWRITE: bool = os.getenv("FORCE_OVERWRITE", "false").lower() == "true"
CHUNKED_SCHEMA_VERSION: str = os.getenv("CHUNKED_SCHEMA_VERSION", "chunked_v1")

# These get assigned in _init_env()
s3 = None
np = None
_pa = None
_pq = None

def _init_env():
    """
    Initialize environment variables and optional libs lazily.
    Safe to call multiple times.
    """
    global S3_BUCKET, S3_RAW_PREFIX, S3_CHUNKED_PREFIX
    global SLIDES_PER_CHUNK, DISABLE_OCR, FORCE_OCR, OCR_BACKEND, PPTX_OCR_STRICT, MIN_IMG_BYTES
    global PARSER_VERSION_PPTX, TOKEN_ENCODER, FORCE_OVERWRITE, CHUNKED_SCHEMA_VERSION
    global s3, np, _pa, _pq

    if getattr(_init_env, "done", False):
        return
    S3_BUCKET = os.getenv("DATA_S3_BUCKET") or S3_BUCKET
    S3_RAW_PREFIX = os.getenv("S3_RAW_PREFIX", S3_RAW_PREFIX).rstrip("/") + "/"
    S3_CHUNKED_PREFIX = os.getenv("S3_CHUNKED_PREFIX", S3_CHUNKED_PREFIX).rstrip("/") + "/"
    SLIDES_PER_CHUNK = int(os.getenv("PPTX_SLIDES_PER_CHUNK", str(SLIDES_PER_CHUNK)))
    DISABLE_OCR = os.getenv("PPTX_DISABLE_OCR", "false").lower() == "true"
    FORCE_OCR = os.getenv("PPTX_FORCE_OCR", "false").lower() == "true"
    OCR_BACKEND = os.getenv("PPTX_OCR_ENGINE", OCR_BACKEND).lower()
    PPTX_OCR_STRICT = os.getenv("PPTX_OCR_STRICT", "false").lower() == "true"
    MIN_IMG_BYTES = int(os.getenv("PPTX_MIN_IMG_SIZE_BYTES", str(MIN_IMG_BYTES)))
    PARSER_VERSION_PPTX = os.getenv("PARSER_VERSION_PPTX", PARSER_VERSION_PPTX)
    TOKEN_ENCODER = os.getenv("TOKEN_ENCODER", TOKEN_ENCODER)
    FORCE_OVERWRITE = os.getenv("FORCE_OVERWRITE", "false").lower() == "true"
    CHUNKED_SCHEMA_VERSION = os.getenv("CHUNKED_SCHEMA_VERSION", CHUNKED_SCHEMA_VERSION)

    # boto3 client
    try:
        import boto3  # local import to avoid hard dependency at import-time
        s3 = boto3.client("s3")
    except Exception as e:
        s3 = None
        log.error("boto3 client could not be created: %s", e)

    # numpy optional
    try:
        import numpy as _np
        np = _np
    except Exception:
        np = None
        log.debug("numpy not available; image handling will be degraded")

    # pyarrow optional (parquet): keep references in module globals
    try:
        import pyarrow as pa
        import pyarrow.parquet as pq
        _pa = pa
        _pq = pq
        log.debug("pyarrow available: %s", getattr(pa, "__version__", "unknown"))
    except Exception:
        _pa = None
        _pq = None
        log.debug("pyarrow not available; parquet operations will fail if used")

    _init_env.done = True

# --- small utilities -------------------------------------------------------
def sha256_hex_bytes(b: bytes) -> str:
    return hashlib.sha256(b).hexdigest()

def sha256_hex_str(s: str) -> str:
    return hashlib.sha256((s or "").encode("utf-8")).hexdigest()

def canonicalize_text(s: str) -> str:
    if not isinstance(s, str):
        s = str(s or "")
    s = unicodedata.normalize("NFKC", s)
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"[ \t]+$", "", ln) for ln in s.split("\n")]
    return "\n".join(lines).strip()

def _load_tiktoken_encoder(name: str):
    try:
        import tiktoken
        enc = None
        try:
            enc = tiktoken.encoding_for_model(name)
        except Exception:
            try:
                enc = tiktoken.get_encoding(name)
            except Exception:
                enc = None
        return enc
    except Exception:
        return None

def _count_tokens(text: str) -> int:
    if not text:
        return 0
    enc = _load_tiktoken_encoder(TOKEN_ENCODER)
    if enc:
        try:
            return len(enc.encode(text))
        except Exception:
            pass
    return len(text.split())

def is_ocr_line_valid(text: str, min_ratio: float = 0.6) -> bool:
    t = (text or "").strip()
    if len(t) < 5:
        return False
    alnum = sum(c.isalnum() for c in t)
    try:
        return (alnum / len(t)) >= min_ratio
    except Exception:
        return False

def dedupe_lines(lines: list) -> list:
    seen, out = set(), []
    for l in lines:
        key = (l or "").strip().lower()
        if key and key not in seen:
            seen.add(key)
            out.append(l)
    return out

# --- s3 helpers ------------------------------------------------------------
def s3_object_exists(key: str) -> bool:
    _init_env()
    if s3 is None:
        raise RuntimeError("boto3 s3 client not available")
    try:
        s3.head_object(Bucket=S3_BUCKET, Key=key)
        return True
    except ClientError as e:
        code = e.response.get("Error", {}).get("Code", "")
        if code in ("404", "NoSuchKey", "NotFound", "NotFoundException"):
            return False
        raise
    except Exception:
        return False

def s3_upload_file_atomic(local_path: str, bucket: str, key: str, content_type: str = "application/octet-stream") -> None:
    _init_env()
    if s3 is None:
        raise RuntimeError("boto3 s3 client not available")
    tmp_key = f"{key}.tmp.{os.getpid()}.{int(time.time())}"
    retries = int(os.getenv("S3_PUT_RETRIES", "3"))
    backoff = float(os.getenv("S3_PUT_BACKOFF", "0.3"))
    for attempt in range(1, retries + 1):
        try:
            s3.upload_file(local_path, bucket, tmp_key, ExtraArgs={"ContentType": content_type})
            copy_source = {"Bucket": bucket, "Key": tmp_key}
            s3.copy_object(CopySource=copy_source, Bucket=bucket, Key=key)
            s3.delete_object(Bucket=bucket, Key=tmp_key)
            return
        except Exception as e:
            log.warning("s3 atomic upload attempt %d failed for %s: %s", attempt, key, e)
            time.sleep(backoff * attempt)
    raise Exception(f"s3 atomic upload failed for {key} after retries")

# --- parquet writer (lazy pyarrow usage) -----------------------------------
class S3ParquetWriter:
    def __init__(self, doc_id: str):
        self.doc_id = doc_id
        self._rows: List[Dict[str, Any]] = []

    def _normalize(self, payload: Dict[str, Any]) -> Dict[str, Any]:
        fields: Dict[str, Any] = {}
        fields["document_id"] = payload.get("document_id") or ""
        fields["file_name"] = payload.get("file_name") or ""
        fields["chunk_id"] = payload.get("chunk_id") or ""
        fields["chunk_type"] = payload.get("chunk_type") or ""
        fields["text"] = payload.get("text") or ""
        try:
            fields["token_count"] = int(payload.get("token_count") or 0)
        except Exception:
            fields["token_count"] = 0
        for k in ("figures", "tags", "layout_tags", "heading_path", "headings"):
            v = payload.get(k, None)
            try:
                fields[k] = json.dumps(v, ensure_ascii=False, sort_keys=True) if v is not None else "[]"
            except Exception:
                fields[k] = "[]"
        fields["file_type"] = payload.get("file_type") or ""
        fields["source_url"] = payload.get("source_url") or ""
        fields["slide_start"] = None
        fields["slide_end"] = None
        if payload.get("slide_range") and isinstance(payload.get("slide_range"), (list, tuple)) and len(payload.get("slide_range")) >= 2:
            try:
                fields["slide_start"] = int(payload["slide_range"][0])
                fields["slide_end"] = int(payload["slide_range"][1])
            except Exception:
                fields["slide_start"] = None
                fields["slide_end"] = None
        fields["timestamp"] = payload.get("timestamp") or ""
        fields["parser_version"] = payload.get("parser_version") or PARSER_VERSION_PPTX
        fields["used_ocr"] = bool(payload.get("used_ocr", False))
        fields["layout"] = payload.get("layout") or ""
        return fields

    def write_payload(self, payload: Dict[str, Any]) -> int:
        self._rows.append(self._normalize(payload))
        return 1

    def finalize_and_upload(self, out_basename: str) -> Tuple[int, str, str, int]:
        _init_env()
        if _pa is None or _pq is None:
            raise RuntimeError("pyarrow is required to write parquet output")
        if not self._rows:
            return 0, "", "", 0
        schema = _pa.schema([
            _pa.field("document_id", _pa.string()),
            _pa.field("file_name", _pa.string()),
            _pa.field("chunk_id", _pa.string()),
            _pa.field("chunk_type", _pa.string()),
            _pa.field("text", _pa.string()),
            _pa.field("token_count", _pa.int64()),
            _pa.field("figures", _pa.string()),
            _pa.field("tags", _pa.string()),
            _pa.field("layout_tags", _pa.string()),
            _pa.field("heading_path", _pa.string()),
            _pa.field("headings", _pa.string()),
            _pa.field("file_type", _pa.string()),
            _pa.field("source_url", _pa.string()),
            _pa.field("slide_start", _pa.int64()),
            _pa.field("slide_end", _pa.int64()),
            _pa.field("timestamp", _pa.string()),
            _pa.field("parser_version", _pa.string()),
            _pa.field("used_ocr", _pa.bool_()),
            _pa.field("layout", _pa.string())
        ])
        cols = {name: [] for name in [f.name for f in schema]}
        for r in self._rows:
            for name in cols:
                cols[name].append(r.get(name) if name in r else None)
        table = _pa.Table.from_pydict(cols, schema=schema)
        existing_md = table.schema.metadata or {}
        new_md = dict(existing_md)
        new_md.update({
            b"schema_version": CHUNKED_SCHEMA_VERSION.encode("utf-8"),
            b"parser_version": PARSER_VERSION_PPTX.encode("utf-8"),
            b"producer": b"pptx_parser",
            b"created_at": datetime.utcnow().isoformat().encode("utf-8")
        })
        table = table.replace_schema_metadata(new_md)
        tmpfile = tempfile.NamedTemporaryFile(mode="wb", delete=False, suffix=".parquet", dir="/tmp")
        tmpfile.close()
        _pq.write_table(table, tmpfile.name, compression="zstd", flavor="spark")
        local_parquet_path = tmpfile.name
        with open(local_parquet_path, "rb") as fh:
            b = fh.read()
        sha = sha256_hex_bytes(b)
        size = os.path.getsize(local_parquet_path)
        parquet_key = out_basename + ".parquet"
        s3_upload_file_atomic(local_parquet_path, S3_BUCKET, S3_CHUNKED_PREFIX + parquet_key, content_type="application/octet-stream")
        try:
            os.unlink(local_parquet_path)
        except Exception:
            pass
        return len(self._rows), S3_CHUNKED_PREFIX + parquet_key, sha, size

# --- small helpers used by parse_file -------------------------------------
def sanitize_payload_for_raw_manifest(doc_id: str, s3_raw_key: str, chunked_s3_key: str, rows: int, sha: str, size: int) -> Dict[str, Any]:
    return {"raw_key": s3_raw_key, "doc_id": doc_id, "chunked_key": chunked_s3_key, "rows": rows, "sha256": sha, "size_bytes": size, "schema_version": CHUNKED_SCHEMA_VERSION, "parser_version": PARSER_VERSION_PPTX, "created_at": datetime.utcnow().isoformat() + "Z"}

def is_valid_table(table: list) -> bool:
    if len(table) < 2 or len(table[0]) < 2:
        return False
    total_cells = sum(len(r) for r in table)
    alpha_cells = sum(1 for row in table for cell in row if any(c.isalpha() for c in (cell or "")))
    return total_cells > 0 and (alpha_cells / total_cells) >= 0.5

def _extract_image_blob_from_shape(shape):
    try:
        img = getattr(shape, "image", None)
        if img and getattr(img, "blob", None):
            return img.blob
    except Exception:
        pass
    try:
        fill = getattr(shape, "fill", None)
        if fill is not None and getattr(fill, "type", None) is not None:
            pic = getattr(fill, "picture", None)
            if pic and getattr(pic, "image", None) and getattr(pic.image, "blob", None):
                return pic.image.blob
    except Exception:
        pass
    return None

def do_ocr(img: Any) -> list:
    """
    img expected as a numpy array (BGR) when using tesseract path.
    This function logs and returns [] on import failures; raising only when strict mode requires it.
    """
    lines = []
    try:
        if OCR_BACKEND == "tesseract":
            try:
                import cv2
                from PIL import Image
                import pytesseract
            except Exception as e:
                log.warning("Tesseract backend requested but imports failed: %s", e)
                if PPTX_OCR_STRICT or FORCE_OCR:
                    raise
                return []
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            _, bin_img = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            raw = pytesseract.image_to_string(Image.fromarray(bin_img), config="--oem 1 --psm 6")
            for l in raw.splitlines():
                if is_ocr_line_valid(l):
                    lines.append(l.strip())
        elif OCR_BACKEND == "rapidocr":
            try:
                from rapidocr_onnxruntime import RapidOCR
            except Exception as e:
                log.warning("RapidOCR backend requested but import failed: %s", e)
                if PPTX_OCR_STRICT or FORCE_OCR:
                    raise
                return []
            ocr = RapidOCR()
            res = ocr(img)
            if res and isinstance(res[0], (list, tuple)):
                for item in res[0]:
                    if len(item) >= 2:
                        text = item[1].strip()
                        if is_ocr_line_valid(text):
                            lines.append(text)
    except Exception:
        return []
    return dedupe_lines(lines)

# --- main parse_file API ---------------------------------------------------
def parse_file(s3_key: str, manifest: dict) -> dict:
    """
    Main entrypoint expected by router.py.
    Validates runtime requirements lazily and returns a dict containing 'saved_chunks'.
    """
    _init_env()

    start_all = time.perf_counter()

    if S3_BUCKET is None:
        raise RuntimeError("S3_BUCKET must be set in environment")

    # Head the object first
    try:
        head_obj = s3.head_object(Bucket=S3_BUCKET, Key=s3_key)
    except Exception as e:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.error("Could not HEAD S3 object %s: %s", s3_key, e)
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True, "error": str(e)}

    content_len = head_obj.get("ContentLength", 0) or 0
    if content_len == 0:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.info("Skipping empty object %s (zero bytes).", s3_key)
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True}

    # Get object
    try:
        raw = s3.get_object(Bucket=S3_BUCKET, Key=s3_key)["Body"].read()
    except Exception as e:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.error("Could not read S3 object %s: %s", s3_key, e)
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True, "error": str(e)}

    # derive doc id
    if isinstance(manifest, dict) and manifest.get("file_hash"):
        doc_id = manifest.get("file_hash")
    else:
        doc_id = sha256_hex_str(raw.decode("latin-1") if isinstance(raw, (bytes, bytearray)) else str(raw))

    out_basename = f"{doc_id}"
    raw_manifest_key = s3_key + ".manifest.json"

    # skip if outputs exist
    if not FORCE_OVERWRITE:
        try:
            if s3_object_exists(raw_manifest_key):
                total_ms = int((time.perf_counter() - start_all) * 1000)
                log.info("Skipping because raw manifest exists: %s", raw_manifest_key)
                return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True}
            if s3_object_exists(S3_CHUNKED_PREFIX + out_basename + ".parquet"):
                total_ms = int((time.perf_counter() - start_all) * 1000)
                log.info("Skipping because parquet chunked file exists: %s", out_basename + ".parquet")
                # attempt to create raw manifest if missing
                try:
                    if not s3_object_exists(raw_manifest_key):
                        head = s3.head_object(Bucket=S3_BUCKET, Key=S3_CHUNKED_PREFIX + out_basename + ".parquet")
                        etag = head.get("ETag", "")
                        if isinstance(etag, str):
                            etag = etag.strip('"')
                        size = head.get("ContentLength", 0)
                        raw_manifest = sanitize_payload_for_raw_manifest(doc_id, s3_key, S3_CHUNKED_PREFIX + out_basename + ".parquet", 0, etag, size)
                        s3.put_object(Bucket=S3_BUCKET, Key=raw_manifest_key, Body=json.dumps(raw_manifest).encode("utf-8"), ContentType="application/json")
                except Exception:
                    pass
                return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True}
        except Exception as e:
            log.warning("Error checking existing outputs for %s: %s", s3_key, e)

    # open pptx
    try:
        from pptx import Presentation
    except Exception as e:
        log.error("pptx import failed: %s", e)
        return {"saved_chunks": 0, "total_parse_duration_ms": 0, "skipped": True, "error": str(e)}

    try:
        prs = Presentation(BytesIO(raw))
    except Exception as e:
        log.error("Failed to open presentation %s: %s", s3_key, e)
        return {"saved_chunks": 0, "total_parse_duration_ms": 0, "skipped": True, "error": str(e)}

    # extract slide content
    slides_content = []
    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        t_slide_start = time.perf_counter()
        text_items, table_items, img_texts = [], [], []
        layout_name = ""
        try:
            layout_name = getattr(getattr(slide, "slide_layout", None), "name", "") or ""
        except Exception:
            layout_name = ""
        for shape in slide.shapes:
            try:
                if getattr(shape, "has_text_frame", False):
                    txt = shape.text or ""
                    if txt.strip():
                        for ln in txt.splitlines():
                            if ln.strip():
                                text_items.append(ln.strip())
                if getattr(shape, "has_table", False):
                    tbl = shape.table
                    rows = []
                    for r in tbl.rows:
                        cols = []
                        for c in r.cells:
                            cols.append((c.text or "").replace("\n", " ").strip())
                        rows.append(cols)
                    norm = [[cell for cell in row] for row in rows]
                    if is_valid_table(norm):
                        header = "| " + " | ".join(norm[0]) + " |"
                        sep = "| " + " | ".join(["---"] * len(norm[0])) + " |"
                        rows_md = ["| " + " | ".join(r) + " |" for r in norm[1:]] if len(norm) > 1 else ["\t".join(r) for r in norm]
                        md_table = "\n".join([header, sep] + rows_md) if len(norm) > 1 else "\n".join(rows_md)
                        table_items.append(md_table)
                blob = _extract_image_blob_from_shape(shape)
                if blob and len(blob) >= MIN_IMG_BYTES:
                    try:
                        from PIL import Image
                        img = Image.open(BytesIO(blob)).convert("RGB")
                        arr = None
                        if np is not None:
                            arr = np.array(img)[:, :, ::-1]
                        else:
                            # No numpy: try to convert to raw bytes and skip OCR if forced off
                            arr = None
                        if arr is not None:
                            ocr_lines = do_ocr(arr)
                            if ocr_lines:
                                img_texts.append("\n".join(ocr_lines))
                    except Exception:
                        # continue; don't fail whole slide for image issues
                        pass
            except Exception:
                continue
        merged_lines = []
        if text_items:
            merged_lines.extend(text_items)
        if table_items:
            merged_lines.extend(table_items)
        if img_texts:
            merged_lines.extend(img_texts)
        merged_lines = [ln for ln in merged_lines if is_ocr_line_valid(ln)]
        merged_lines = dedupe_lines(merged_lines)
        slide_parse_ms = (time.perf_counter() - t_slide_start) * 1000.0
        slides_content.append({
            "slide_number": slide_num,
            "raw_lines": merged_lines,
            "has_text": bool(text_items),
            "has_images_text": bool(img_texts),
            "tables": table_items,
            "parse_duration_ms": slide_parse_ms,
            "layout": layout_name or ""
        })

    # build and write chunks
    saved = 0
    total_slides = len(slides_content)
    writer = S3ParquetWriter(doc_id=doc_id)
    try:
        def _sanitize_payload_for_weaviate(payload: dict) -> dict:
            range_keys = {"row_range", "slide_range", "token_range", "audio_range", "line_range"}
            for k in list(payload.keys()):
                v = payload.get(k)
                if k in range_keys and isinstance(v, (list, tuple)):
                    try:
                        if len(v) == 2 and all(isinstance(x, (int, str)) for x in v):
                            payload[k] = [int(v[0]), int(v[1])]
                        else:
                            payload[k] = json.dumps(v)
                    except Exception:
                        payload[k] = json.dumps(v)
            if "headings" in payload and isinstance(payload["headings"], (list, tuple)):
                payload["headings"] = [str(x) for x in payload["headings"]]
            if "heading_path" in payload and isinstance(payload["heading_path"], (list, tuple)):
                payload["heading_path"] = [str(x) for x in payload["heading_path"]]
            if "tags" in payload and isinstance(payload["tags"], (list, tuple)):
                payload["tags"] = [str(x) for x in payload["tags"]]
            return payload

        for i in range(0, total_slides, SLIDES_PER_CHUNK):
            chunk_slides = slides_content[i:i + SLIDES_PER_CHUNK]
            start = chunk_slides[0]["slide_number"]
            end = chunk_slides[-1]["slide_number"]
            chunk_id = f"{doc_id}_slides_{start}_{end}"
            t_chunk_start = time.perf_counter()
            merged, used_ocr, slides_sum_ms, layouts = [], False, 0.0, []
            for slide in chunk_slides:
                merged.append(f"## Slide {slide['slide_number']}")
                for ln in slide["raw_lines"]:
                    merged.append(ln)
                if slide["tables"]:
                    merged.extend(slide["tables"])
                if slide["has_images_text"]:
                    used_ocr = True
                if not slide["has_text"] and slide["has_images_text"]:
                    used_ocr = True
                layouts.append(str(slide.get("layout", "") or ""))
                slides_sum_ms += float(slide.get("parse_duration_ms", 0.0))
            clean = [ln for ln in merged if is_ocr_line_valid(ln)]
            clean = dedupe_lines(clean)
            final_text = "\n\n".join(clean)
            token_count = _count_tokens(final_text)
            merge_write_ms = (time.perf_counter() - t_chunk_start) * 1000.0
            duration_ms = int(slides_sum_ms + merge_write_ms)
            dedup_layouts = []
            for l in layouts:
                if l and l not in dedup_layouts:
                    dedup_layouts.append(l)
            layout_str = ";".join(dedup_layouts) if dedup_layouts else ""
            payload = {
                "document_id": doc_id or "",
                "file_name": os.path.basename(s3_key),
                "chunk_id": chunk_id or "",
                "chunk_type": "slides",
                "text": final_text or "",
                "token_count": int(token_count or 0),
                "figures": "[]",
                "embedding": None,
                "file_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "source_url": f"s3://{S3_BUCKET}/{s3_key}",
                "slide_range": [int(start), int(end)],
                "timestamp": (datetime.utcfromtimestamp(int(os.getenv('SOURCE_DATE_EPOCH'))).isoformat() + "Z") if os.getenv("SOURCE_DATE_EPOCH") else (datetime.utcnow().isoformat() + "Z"),
                "parser_version": PARSER_VERSION_PPTX,
                "tags": manifest.get("tags", []) if isinstance(manifest, dict) else [],
                "layout_tags": None,
                "layout": layout_str,
                "used_ocr": bool(used_ocr),
                "heading_path": [],
                "headings": [],
                "line_range": None
            }
            payload = _sanitize_payload_for_weaviate(payload)
            writer.write_payload(payload)
            log.info("Buffered slides %d-%d (tokens=%d)", start, end, token_count)
            saved += 1
    except Exception as e:
        log.exception("Fatal error while buffering chunks for %s: %s", s3_key, str(e))
        return {"saved_chunks": 0, "total_parse_duration_ms": int((time.perf_counter() - start_all) * 1000), "skipped": True, "error": str(e)}

    # finalize and upload parquet
    try:
        if saved == 0:
            total_ms = int((time.perf_counter() - start_all) * 1000)
            log.info("No chunks produced for %s", s3_key)
            return {"saved_chunks": 0, "total_parse_duration_ms": total_ms}
        count, uploaded_s3_key, sha, size = writer.finalize_and_upload(out_basename)
        total_ms = int((time.perf_counter() - start_all) * 1000)
        try:
            raw_manifest = sanitize_payload_for_raw_manifest(doc_id, s3_key, uploaded_s3_key, count, sha, size)
            s3.put_object(Bucket=S3_BUCKET, Key=raw_manifest_key, Body=json.dumps(raw_manifest).encode("utf-8"), ContentType="application/json")
        except Exception:
            log.warning("Failed to write raw manifest for %s", s3_key)
        log.info("Wrote %d chunks for %s → %s (%d ms)", count, s3_key, uploaded_s3_key, total_ms)
        return {"saved_chunks": count, "total_parse_duration_ms": total_ms, "skipped": False}
    except Exception as e_up:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.error("Failed to upload chunked file for %s error=%s", s3_key, str(e_up))
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True, "error": str(e_up)}

# --- CLI runner (preserve behavior) ---------------------------------------
if __name__ == "__main__":
    _init_env()
    log.info("Starting pptx -> parquet parser")
    if S3_BUCKET is None or s3 is None:
        log.error("S3_BUCKET or s3 client not configured; aborting CLI run")
        sys.exit(1)
    paginator = s3.get_paginator("list_objects_v2")
    for page in paginator.paginate(Bucket=S3_BUCKET, Prefix=S3_RAW_PREFIX):
        for obj in page.get("Contents", []):
            key = obj["Key"]
            if not key.lower().endswith(".pptx"):
                continue
            log.info("Routing parse_file for s3://%s/%s", S3_BUCKET, key)
            manifest_key = key + ".manifest.json"
            try:
                mf = s3.get_object(Bucket=S3_BUCKET, Key=manifest_key)
                manifest = json.load(mf["Body"])
            except ClientError:
                manifest = {}
            except Exception:
                manifest = {}
            try:
                result = parse_file(key, manifest)
                log.info("Result for %s: %s", key, result)
            except Exception as e:
                log.exception("Failed to parse %s: %s", key, e)
